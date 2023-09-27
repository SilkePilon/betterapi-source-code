'''
The following code will create a pool of workers of 4 processes with 4 threads under 
each process. They rely on redis for synchronization, so you can run other instances 
as you like without worrying about creating conflicts.
'''
from fastapi_queue import QueueWorker
from loguru import logger
import asyncio  
import aioredis
from os import environ
from Bard import Chatbot
from instabot import Bot
import signal 
import sys 
import os 
# main.py
os.system("killall chrome")
import asyncio
import collections.abc
from pptx import Presentation
from pptx.util import Inches
import random 
import re
import requests
import json
import os
import platform
import re
import sys
import time
from importlib import metadata as importlib_metadata
import markdownify
import undetected_chromedriver as uc
import urllib3
from colorama import Fore
from pyvirtualdisplay import Display
from selenium.common import exceptions as SeleniumExceptions
from selenium.webdriver.common.action_chains import ActionChains
from selenium.webdriver.common.by import By
from selenium.webdriver.common.keys import Keys
from selenium.webdriver.support import expected_conditions as EC
from selenium.webdriver.support.ui import WebDriverWait

urllib3.disable_warnings()

from fastapi import FastAPI
import asyncio
import json
import os
import platform
import re
import subprocess
import time
import urllib.parse
import uvicorn
import cloudscraper
import markdownify
import undetected_chromedriver as uc
import urllib3
from pyvirtualdisplay import Display
from ratelimit import limits
from selenium.common import exceptions as SeleniumExceptions
from selenium.webdriver.common.action_chains import ActionChains
from selenium.webdriver.common.by import By
from selenium.webdriver.common.keys import Keys
from selenium.webdriver.support import expected_conditions as EC
from selenium.webdriver.support.ui import WebDriverWait
from gtts import gTTS
from slowapi import Limiter, _rate_limit_exceeded_handler
from slowapi.util import get_remote_address
from slowapi.errors import RateLimitExceeded
from starlette.requests import Request
from fastapi.responses import FileResponse
urllib3.disable_warnings()
from io import BytesIO
import string
from api_analytics.fastapi import Analytics
import traceback
import requests
import random
from fastapi import Response
from fastapi.encoders import jsonable_encoder
from fastapi.responses import JSONResponse
from youdotcom import Webdriver
from starlette.responses import PlainTextResponse, RedirectResponse
from apscheduler.schedulers.background import BackgroundScheduler
from apscheduler.schedulers.blocking import BlockingScheduler
from selenium.webdriver.common.by import By
from selenium.webdriver.support.ui import WebDriverWait
from selenium.webdriver.support import expected_conditions as EC
from scrapingbee import ScrapingBeeClient
try:
    display = Display()
except FileNotFoundError as e:
    if "No such file or directory: 'Xvfb'" in str(e):
        raise ValueError("Headless machine detected. Please install Xvfb to start a virtual display: sudo apt install xvfb")
    raise e
# display.start()

options = uc.ChromeOptions()
options.add_argument("--headless")
options.add_argument("--start_maximized")
# options.add_argument("--disable-extensions")
# options.add_argument("--disable-application-cache")
# options.add_argument("--disable-gpu")
# options.add_argument("--no-sandbox")
# options.add_argument("--disable-setuid-sandbox")
# options.add_argument("--disable-dev-shm-usage")
options.add_argument('--user-agent=Mozilla/5.0 (Windows NT 10.0; Win64; x64; rv:109.0) Gecko/20100101 Firefox/109.0')
# options.add_experimental_option("excludeSwitches", ["enable-automation"])
# options.add_experimental_option('useAutomationExtension', False)
options.add_argument('--disable-blink-features=AutomationControlled')
driver = uc.Chrome(options=options, driver_executable_path=f"/usr/lib/chromium-browser/chromedriver", version_main=110, use_subprocess=True)



# bot = Bot()
# bot.login(username="BetterAPI", password="Landrover@01")



# #  ---------- EDIT ----------
# email = 'silke2007minecraft@gmail.com\n' # replace email
# password = 'Landrover01\n' # replace password
# #  ---------- EDIT ----------


# # driver = Webdriver(webdriver_path="/usr/lib/chromium-browser/chromedriver", hide=True, headless=True).driver
# wait = WebDriverWait(driver, 40)
# url = 'https://accounts.google.com/ServiceLogin?service=accountsettings&continue=https://myaccount.google.com%3Futm_source%3Daccount-marketing-page%26utm_medium%3Dgo-to-account-button'
# driver.get(url)


# wait.until(EC.visibility_of_element_located((By.NAME, 'identifier'))).send_keys(email)
# driver.save_screenshot("why.png")
# try:
#     wait.until(EC.visibility_of_element_located((By.NAME, 'password'))).send_keys(password)
# except:
#     driver.save_screenshot("why2.png")
# print("You're in!! enjoy")




from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles
from datetime import datetime
from PIL import Image
from typing import Optional, Any
from fastapi import FastAPI, Request
from fastapi.responses import JSONResponse
from fastapi_queue import DistributedTaskApplyManager
import aioredis
from fastapi import FastAPI, Request, status
from fastapi.encoders import jsonable_encoder
from fastapi.exceptions import RequestValidationError
from fastapi.responses import JSONResponse
from pydantic import BaseModel
import urllib.parse
import base64
import cfscrape
import paramiko
from urllib.parse import unquote, quote_plus ,quote
import urllib.parse
queueworker = None
from random_user_agent.user_agent import UserAgent
from random_user_agent.params import SoftwareName, OperatingSystem
software_names = [SoftwareName.CHROME.value]
operating_systems = [OperatingSystem.WINDOWS.value, OperatingSystem.LINUX.value]
user_agent_rotator = UserAgent(software_names=software_names, operating_systems=operating_systems, limit=100)

# client = ScrapingBeeClient(api_key='ONSH6064PINTZCDC0VCH6SOBAAUNH374A85GJQHYJ9S9WBGBDP97BSWJD6JMS963SLU7VOVJXMG5V6ZO')
    

# browser={"browser": "chrome", "platform": "android", "mobile": True,'desktop': False}



ssh = paramiko.SSHClient() 
ssh.load_host_keys(os.path.expanduser(os.path.join("~", ".ssh", "known_hosts")))

user_agent = "Mozilla/5.0 (Windows NT 10.0; Win64; x64; rv:109.0) Gecko/20100101 Firefox/109.0"


def setnewcookies():
    print("-- GETTING NEW COOKIES --")
    ssh.connect("192.168.2.151", username="silke", password="Landrover01")
    # user_agent_random = user_agent_rotator.get_random_user_agent()
    # print(user_agent_random)
    # stidn2d, dwa, dwawd =ssh.exec_command("rm /home/silke/CF-Clearance-Scraper/cookies.json")
    stdin, stdout, stderr = ssh.exec_command(f"python3 /home/silke/CF-Clearance-Scraper/main.py -u https://you.com/api/streamingSearch")
    # exit_status = stdout.channel.recv_exit_status()
    
    cf_clearence = str(stdout.read().decode("utf-8")).replace("\n", "")
            
    
    
    global cf_bypass
    cf_bypass = cf_clearence
    global user_agent
   
    print(cf_clearence)
    print(user_agent)
    global scraper
    scraper = cloudscraper.CloudScraper(debug=False)
    scraper.headers.update({'User-Agent': f'{user_agent}'})
    global cookiesvar
    cookiesvar = {'cf_clearance': f'{cf_bypass}'}
    ssh.close()
    print("-- COOKIES SET --")


# setnewcookies()
# scheduler = BackgroundScheduler()
# scheduler.add_job(setnewcookies, 'interval', minutes=7)
# scheduler.start()




headerslist = {
    'Content-Type': 'application/json',
}

sessuion_list = {
    'cmd': 'sessions.list',
    'maxTimeout': 60000,
}

response = requests.post('http://192.168.2.151:8191/v1', headers=headerslist, json=sessuion_list).json()
if not "MAIN" in response['sessions']:
    create_session = {
        'cmd': 'sessions.create',
        'session': 'MAIN',
        'maxTimeout': 60000,
    }

    response = requests.post('http://192.168.2.151:8191/v1', headers=headerslist, json=create_session).json()
    print(response)




def sendmessagetochat(redis,mysql, inputs, contextid, ip, url,key, powerpointer, parameters):
    try:
        # with open("log.json",'r+') as file:
        #     # First we load existing data into a dict.
        #     file_data = json.load(file)
        #     # Join new_data with file_data inside emp_details
        
        #     file_data["users"][f"{key}"]["chat_history"].append(message)
        #     # Sets file's current position at offset.
        #     file.seek(0)
        #     # convert back to json.
        #     json.dump(file_data, file, indent = 4)
        # message = urllib.parse.quote(message)
        
        message = inputs
        
        if powerpointer == True or powerpointer == "True":
            message = f"""Write a presentation slash powerpoint text about the user's topic. You only answer with the presentation text. Follow the structure of the example.
Notice
-You do all the presentation text for the user.
-You write the texts no longer than 250 characters!
-You make very short titles!
-You make the presentation easy to understand.
-The presentation has a table of contents.
-The presentation has a summary.
-At least 8 slides.
-You shall write the presentation text in the language of the user
-Check if the user want diffrent infomation per slide in the presentation if so do that

Example! - Stick to this formatting exactly!
#Title: TITLE OF THE PRESENTATION

#Slide: 1
#Header: table of contents
#Content: 1. CONTENT OF THIS POWERPOINT
2. CONTENTS OF THIS POWERPOINT
3. CONTENT OF THIS POWERPOINT
...

#Slide: 2
#Header: TITLE OF SLIDE
#Content: CONTENT OF THE SLIDE

#Slide: 3
#Header: TITLE OF SLIDE
#Content: CONTENT OF THE SLIDE

#Slide: 4
#Header: TITLE OF SLIDE
#Content: CONTENT OF THE SLIDE

#Slide: 5
#Headers: summary
#Content: CONTENT OF THE SUMMARY

#Slide: END

The user wants a presentation about or wants part of the presentation about: {message}
"""
        
        
        start = time.time()
        CloudflareChallengeError = False
        typeof = ""
        status_code = 200
        if contextid == "" or None:
            contextid = ""
        print(message)
        message = quote(message)
        
        global chat
        
        chat = []
        # print(user_agent)
        # headers = {
        #     # "user-agent": f"{user_agent}",
        #     "Accept": "text/event-stream",
        #     "Connection": "keep-alive",
        #     "Sec-Fetch-Mode": "cors",
        #     "Sec-Fetch-Site": "same-origin",
        #     "Sec-GPC": "1",
        #     "Referer": "https://you.com/search?q=hello&fromSearchBar=true&tbm=youchat",
        #     "Cookie": f"uuid_guest=dummystring;cf_clearance={cf_bypass}",
        # }
        payload = {
            "q": message,
            "chat": "",
            "queryTraceId": "",
            "domain": "youchat",
            "page": "1",
            "count": "10",
            "safeSearch": "Off",
            "onShoppingPage": "false",
            "freshness": "Month",
            "mkt": "",
            "responseFilter": "WebPages,Translations,TimeZone,Computation,RelatedSearches",
            "sharedChatId": f"{contextid}"
        }
        output = ""
        try:
            headers = {
                'Content-Type': 'application/json',
            }


            json_data = {
                'cmd': 'request.get',
                'session': "MAIN",
                'url': f'https://you.com/api/streamingSearch?q={message}&page=1&count=10&safeSearch=Moderate&onShoppingPage=false&mkt=&responseFilter=WebPages,Translations,TimeZone,Computation,RelatedSearches&domain=youchat&queryTraceId=&chat=&chatId=',
                'maxTimeout': 60000,
            }


            # session = requests.post('http://192.168.2.151:8191/v1', headers=headers, json=create).json()
            # print(session)
            response = requests.post('http://192.168.2.151:8191/v1', headers=headers, json=json_data).json()
            try:
                response = str(response['solution']['response']).replace('<html><head><meta name="color-scheme" content="light dark"></head><body><pre style="word-wrap: break-word; white-space: pre-wrap;">', '').split("\n")
            except:
                print(response)
            
            
            # print(response)
                 
            # response = scraper.get(f"https://you.com/api/streamingSearch", params=payload, headers=headers, cookies=cookiesvar, stream=True)
            # # ad = scraper.get(f"https://you.com/api/ads?q={message}", headers=headers, cookies=cookiesvar).json()['suggestions']
            CloudflareChallengeError = False
            
            
        except Exception as e:
            print(traceback.format_exc())
            # setnewcookies()
            CloudflareChallengeError = True
        
        if CloudflareChallengeError == True:
            output = "cloudflare error"
            status_code = 503
            
        if CloudflareChallengeError == False:
            
            for line in response:
                if line:
                    try:
                        key, value = line.split(":", 1)
                        key = key.strip()
                        value = value.strip()
                        if key == "data":
                            if value == "I'm Mr. Meeseeks. Look at me.":
                                break
                            data = json.loads(value)
                            if "youChatToken" in data:
                                output += data["youChatToken"]
                            # if "youChatAppData" in data:
                            #     apps = data``
                    except Exception as e:
                        # print(traceback.format_exc())
                        pass
        if len(chat) > 4:
            chat = chat[:-4]

        # print(chat)

        out = str(re.sub(r"\[.+?\]\(.+?\)", "", output)).lstrip()
        # out = out[1:]
        # msg = markdownify.markdownify(out)
        
        if out == "":
            out = "error, please vote and read: https://discord.com/channels/958041025539080202/1099742780563128452"

        # subprocess.call(["taskkill", "/im", "chromedriver.exe"],shell=True)
        # subprocess.call(["pkill", "-f", "chromedriver"], shell=True)
        # subprocess.call(["killall", "-m", "chromedriver"], shell=True)
        
        
        
        # ad = requests.get(f"https://you.com/api/ads?q={message}").json()
        # if not len(ad['suggestions']) == 0 or ad['suggestions' == []]:
        #     out = out + f"\nAD: {ad['suggestions']}"
            
        
        
        timedate = time.time() - start
        timedate = time.strftime("%S", time.gmtime(timedate))
        # out = base64.b64encode(bytes('your string', 'utf-8')) # bytes
        
        
        
        
        return {
  "details": {
    "best_of_sequences": [
      {
        "finish_reason": "length",
        "generated_text": f"{str(out)}"
      }
    ],
    "finish_reason": "length",
    "generated_tokens": len(str(out))
  },
  "generated_text": f"{str(out)}",
  "time": str(timedate),
  "status_code": status_code
}
    except Exception as e: 
        print(traceback.format_exc())
        return {
  "error": "Request failed during generation"
}
    
    
    
    
    
    
    
    
    
    
    
    
    
def imageAPI(redis,mysql, prompt, ip, url,key, style):
    if type(style) == str:
        return {"url": "please use numbers for the style","stages": "", "time": str(0), 'status_code': 500}
    try:
        start = time.time()
        status_code = 200
        with open('token.txt', 'r') as F:
            auth = F.read()
        
        
        
            
        print(auth)
        
        headers = {
            'authority': 'paint.api.wombo.ai',
            'accept': '*/*',
            'accept-language': 'nl-NL,nl;q=0.9,en;q=0.8',
            'authorization': str(auth),
            'content-type': 'text/plain;charset=UTF-8',
            'dnt': '1',
            'origin': 'https://dream.ai',
            'referer': 'https://dream.ai/',
            'sec-ch-ua': '"Chromium";v="112", "Google Chrome";v="112", "Not:A-Brand";v="99"',
            'sec-ch-ua-mobile': '?0',
            'sec-ch-ua-platform': '"Windows"',
            'sec-fetch-dest': 'empty',
            'sec-fetch-mode': 'cors',
            'sec-fetch-site': 'cross-site',
            'user-agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/112.0.0.0 Safari/537.36',
            'x-app-version': 'WEB-1.90.1',
        }

        data = '{"is_premium":false,"input_spec":{"prompt":"' + prompt + '","style":'+ style +',"display_freq":10}}'

        response = requests.post('https://paint.api.wombo.ai/api/v2/tasks', headers=headers, data=data).json()
        print(response)
        max_loops = 5
        if "Token expired" or "Invalid token" in str(response):
            new_token = requests.get("http://192.168.2.63:5000/get_token",timeout=10000).json()
            new_token = new_token['token']
            print(new_token)
            with open('token.txt', 'w') as X:
                 X.write(new_token)
            new = {
                'authority': 'paint.api.wombo.ai',
                'accept': '*/*',
                'accept-language': 'nl-NL,nl;q=0.9,en;q=0.8',
                'authorization': str(new_token),
                'content-type': 'text/plain;charset=UTF-8',
                'dnt': '1',
                'origin': 'https://dream.ai',
                'referer': 'https://dream.ai/',
                'sec-ch-ua': '"Chromium";v="112", "Google Chrome";v="112", "Not:A-Brand";v="99"',
                'sec-ch-ua-mobile': '?0',
                'sec-ch-ua-platform': '"Windows"',
                'sec-fetch-dest': 'empty',
                'sec-fetch-mode': 'cors',
                'sec-fetch-site': 'cross-site',
                'user-agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/112.0.0.0 Safari/537.36',
                'x-app-version': 'WEB-1.90.1',
            }

            data = '{"is_premium":false,"input_spec":{"prompt":"' + prompt + '","style":'+ style +',"display_freq":10}}'

            new_R = requests.post('https://paint.api.wombo.ai/api/v2/tasks', headers=new, data=data).json()
            print(new_R)
            id = new_R['id']
        else:
            id = response['id']
        
            
           
            
        
        while True:
            if max_loops > 5:
                break
            headers = {
                'authority': 'paint.api.wombo.ai',
                'accept': '*/*',
                'accept-language': 'nl-NL,nl;q=0.9,en;q=0.8',
                'authorization': 'bearer eyJhbGciOiJSUzI1NiIsImtpZCI6ImU3OTMwMjdkYWI0YzcwNmQ2ODg0NGI4MDk2ZTBlYzQzMjYyMjIwMDAiLCJ0eXAiOiJKV1QifQ.eyJwcm92aWRlcl9pZCI6ImFub255bW91cyIsImlzcyI6Imh0dHBzOi8vc2VjdXJldG9rZW4uZ29vZ2xlLmNvbS9wYWludC1wcm9kIiwiYXVkIjoicGFpbnQtcHJvZCIsImF1dGhfdGltZSI6MTY4MjQ5ODM2NCwidXNlcl9pZCI6Inh4b0RqYWpkaGloZjRyYWFqNzRPYVZJb3VRODMiLCJzdWIiOiJ4eG9EamFqZGhpaGY0cmFhajc0T2FWSW91UTgzIiwiaWF0IjoxNjgyNDk4MzY0LCJleHAiOjE2ODI1MDE5NjQsImZpcmViYXNlIjp7ImlkZW50aXRpZXMiOnt9LCJzaWduX2luX3Byb3ZpZGVyIjoiYW5vbnltb3VzIn19.KULREbpzsN1NgSmJyGai_t0v14PViZ1OfP4Ok8CbVa-XHN8P9MquC79p7Z3oRtuwLjC-aQyIiZCnoyARpviPP2DOntjvNyV3Ld6_U81eEKWr_k-1V4hQIH9r17uEFAFg1_Gw4p5mm2cwN_KkwQt87IZMtWSxWtH7_6vPRIGI_HUTtpiI2KQb8Oekl5uC5b0tG1_Y_qKaPP94mdUiLByEL4YClznvJkDBUGZjR_EHLXX4oB5FQvNSewsM4xXbXGlCw9ouQqjFrBWTEvR2fL-UQHOMwy0Un5nP7nzEK4dn0FTzpiE8JrvT9Hn_3-IJYIYoUFTv0MIs-558WVobkWKyzg',
                'dnt': '1',
                'origin': 'https://dream.ai',
                'referer': 'https://dream.ai/',
                'sec-ch-ua': '"Chromium";v="112", "Google Chrome";v="112", "Not:A-Brand";v="99"',
                'sec-ch-ua-mobile': '?0',
                'sec-ch-ua-platform': '"Windows"',
                'sec-fetch-dest': 'empty',
                'sec-fetch-mode': 'cors',
                'sec-fetch-site': 'cross-site',
                'user-agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/112.0.0.0 Safari/537.36',
                'x-app-version': 'WEB-1.90.1',
            }

            response = requests.get(f'https://paint.api.wombo.ai/api/v2/tasks/{id}', headers=headers).json()
            if response['state'] == "completed":
                stages = str(response['photo_url_list']).replace("images.wombo.art", "api.betterapi.net").replace('"', "")
                final = str(response['result']['final']).replace("images.wombo.art", "api.betterapi.net")
                break
            time.sleep(4)
        
        
        timedate = time.time() - start
        timedate = time.strftime("%S", time.gmtime(timedate))
        # path = str(final).replace("api.betterapi.net", "images.wombo.art")
        # with open(f'upload.png', 'wb') as handle:
        #     response = requests.get(path, stream=True)

        #     if not response.ok:
        #         print(response)

        #     for block in response.iter_content(1024):
        #         if not block:
        #             break

        #         handle.write(block)
        
        # bot.upload_photo(file="upload.png", caption="#ai")

        
        return {"url": str(final),"stages": stages, "time": str(timedate), 'status_code': status_code}
    except Exception as e: 
        print(traceback.format_exc())
        return {"error": "Error, not likly made by you: "+str(e)}
    
    
    
    
    
    
    
    
    
    
    
    
    
    
    
    
    
        
        
def getapps(redis,mysql, message, app_name, ip, url,key):
    try:
        # with open("log.json",'r+') as file:
        #     # First we load existing data into a dict.
        #     file_data = json.load(file)
        #     # Join new_data with file_data inside emp_details
        
        #     file_data["users"][f"{key}"]["chat_history"].append(message)
        #     # Sets file's current position at offset.
        #     file.seek(0)
        #     # convert back to json.
        #     json.dump(file_data, file, indent = 4)
        message = urllib.parse.quote(message)
        
        start = time.time()
        CloudflareChallengeError = False
        typeof = ""
        status_code = 200
        
        
        global chat
        chat = []
        headers = {
            "user-agent": f"{user_agent}",
            "Accept": "text/event-stream",
            "Connection": "keep-alive",
            "Sec-Fetch-Mode": "cors",
            "Sec-Fetch-Site": "same-origin",
            "Sec-GPC": "1",
            "Referer": "https://you.com/search?q=hello&fromSearchBar=true&tbm=youchat",
            "Cookie": f"uuid_guest=dummystring;cf_clearance={cf_bypass}",
        }
        payload = {
            "q": message,
            "chat": "",
            "page": "1",
            "count": "10",
            "safeSearch": "Off",
            "onShoppingPage": "false",
            "freshness": "Month",
            "mkt": "",
            "responseFilter": "WebPages,Translations,TimeZone,Computation,RelatedSearches",
            "appName": f"{app_name}"
        }
        try:
            response = scraper.get(f"https://you.com/api/streamingSearch", params=payload, headers=headers,cookies=cookiesvar, stream=True)
            CloudflareChallengeError = False
            typeof = "api"
            
            
        except cloudscraper.exceptions.CloudflareChallengeError as e:
            print(e)
            setnewcookies()
            response = scraper.get(f"https://you.com/api/streamingSearch", params=payload, headers=headers,cookies=cookiesvar, stream=True)
            CloudflareChallengeError = True
            typeof = "api"
        
        
        if CloudflareChallengeError == True:
           
            # for line in response:
            #     if line:
            #         decoded_line = str(line)
            #         key, value = decoded_line.split(":", 1)
            #         key = key.strip()
            #         value = value.strip()
            #         if key == "data":
            #             if value == "I'm Mr. Meeseeks. Look at me.":
            #                 break
            #             if value == "undefined":
            #                 output = "😔 Due to high demand, I'm experiencing issues briefly. Please try again later or use the All tab to get an answer in the meantime."
            #                 break
            #             data = json.loads(value)
            #             if "ydcAppName" in data:
            #                 output = data
            output = "cloudflare error"
            status_code = 503
        if CloudflareChallengeError == False:
            
            for line in response.iter_lines():
                if line:
                    decoded_line = line.decode("utf-8")
                    key, value = decoded_line.split(":", 1)
                    key = key.strip()
                    value = value.strip()
                    if key == "data":
                        if value == "I'm Mr. Meeseeks. Look at me.":
                            break
                        if value == "undefined":
                            output = "😔 Due to high demand, I'm experiencing issues briefly. Please try again later or use the All tab to get an answer in the meantime."
                            status_code = 503
                            break
                        data = json.loads(value)
                        if "ydcAppName" in data:
                            output = data
        if len(chat) > 4:
            chat = chat[:-4]

        # print(chat)

        # out = str(re.sub(r"\[.+?\]\(.+?\)", "", output[1:])).lstrip()
        # out = out[1:]
        # msg = markdownify.markdownify(out)

        # subprocess.call(["taskkill", "/im", "chromedriver.exe"],shell=True)
        # subprocess.call(["pkill", "-f", "chromedriver"], shell=True)
        # subprocess.call(["killall", "-m", "chromedriver"], shell=True)
        timedate = time.time() - start
        timedate = time.strftime("%S", time.gmtime(timedate))
        # out = base64.b64encode(bytes('your string', 'utf-8')) # bytes
        return {"appdata": output, "time": str(timedate), "status_code": status_code}
    except Exception as e: 
        print(e)
        return {"error": "Error, not likly made by you: "+str(e)}
        
def stable_diffusion(redis,mysql, message, ip, url, model):
    try:
        start = time.time()
        CloudflareChallengeError = False
        typeof = ""
        
        if model == "openjourney":
            model = "stable-diffusion-v2.1"
        
        
        global chat
        chat = []
        headers = {
            "Accept": "text/event-stream",
            "Connection": "keep-alive",
            "Sec-Fetch-Mode": "cors",
            "Sec-Fetch-Site": "same-origin",
            "Sec-GPC": "1",
            "Referer": "https://you.com/search?q=hello&fromSearchBar=false&tbm=youchat",
            "Cookie": f"uuid_guest=04147a25-a6fd-4957-96e7-d8113af17df7;appSession=eyJhbGciOiJkaXIiLCJlbmMiOiJBMjU2R0NNIiwiaWF0IjoxNjc1NDU0ODU0LCJ1YXQiOjE2NzY4MjI1MTEsImV4cCI6MTY3OTQxNDUxMX0..9jbxFr1bxjlCDOeM.Ak2klxC_AfwgGjLrgURAEMj-0QJAz0Ia7fQvL4yJw-dKO1CUYmKG6uz07S6td711B5vL_1ZFIbAo-CMfZEpDOHGryajC34gX8Zlg1WY-clqPdkeXGLQ0Tw4QFnIY-SQuSMZCe1fHr6lDzwkjwSi20RLKIL4QNcdu-g4-hYEQlj62lxY85pJgO75ZOPjVifx6EnlBa80DjjIKBTdHXO42O-1iyCrdLC1SSVk9vzfJBqaGumMjPNYholjrye8C8LQZLXfI6ikGplCkOyCjKO9CdF0BqE7jbu6tm_AACSWMuALeftu9GRqA6cJt-xE4xGfGbx71LtBynzLeBHvVKjAzaTgLJPBjdC99CX6enNNUGDGJqJZ5EhUWJfMwnEtD7tjSB3y-n0u2E3QSbsx5cKfSRazpIwnjmeZTkMvayGkeTr5W-bbfMwCafrByTrGBj0972Bv92n8grqs9hUjp-uDHaH_PJnsTfoqOZQUavTM1qbxrLR6W-T1lHcpvM6i5pHXMm5UOh3xhr6nyOhbFBaAqX_UZz0w4jjpgMW9w-K5JR79IzGTbg0VoFfvCq6zPn4AykQz3ZBhihLo17bqTgU8LMkxnG4E66a_RJszFbrlUPSn_mdmGJkcpcTl43Hr-iCuqeBinv2pasPZosjOxnoMl59BRkmnOBlJ-D5uUFF375n81pTH7snzHCBm-fPc6DuY_iLol7SIAuKnQV9K46b9oj0c0CsVK1erhNUkO0B4-OmPLBRD3aDW4IWXgqnZ2nYtQqJljx60Cn8BN-vMHLZuOrG2p4F561FOYoO3inR9DYk0QDgRMD-6J2VLVAX5g0m0fBsJunO-rmtD5DRtKJPgFlOw57Fw9I7Bv7FKtbYILYqH7Z4kaEofmiIE1jHbETDaeSEQtrnMlVI2477S7bl97oyXHMEEsor_kHnxKN0WV0UXUlOnvokmkabzzsW208bnFCE3XM6nUWm6sdKlC9GGks9ild0mG4KQnWNAqUElE-tfhOwiD2yh6etYkQT6ZK15XDaTB45SCk8cfoMRm6Nl-Iyd5DMmbFlrcjJyVt2c7ez8fAHfl46De8-6f_oljQDlv2Nqd5Z5iiBNiefkQ2e2TpnWZMn1p6zaShcxXofFUdQXeNK4kNVjRPqZxuoxC-7-Rxv2E57L-YATOhamff_aGb8PuZW4XTgr3m5V8-_JtqSKNf8ce659BP9udDSqvDbrPhBPnC2VbJjagYzJUfRxkGLJppOYtWzrUUOVauOaDHjZpsIfXxQJGBMegwTXWIsDdzukWGdEwY08YyNOwlmZ3-MUNqXQlX6fCEiEbQJldfhZj77Rx6N1OHS19cJvqarJ7ARwFzlnDm3ew2hJOu4UzqKtfumPqPryIR0-wCyEHwAJ7-eOngnmGwatdyaWiXS0f7KgaDQJ9WBobjyhqIFlndNG7snnkA4SLWy_qjmjpN2-5Tu_uBQR7fhs9kNU8uAVDqlfE1S2O6arSRW9SD4893QofzouC-awmIfAroAhegQzDuVcj7RrCEsp2ujxK8s1pPg40wgnR2_mRDzC_wqnb3OYU2N0VhpqDaJUceScpkWmVE73rgKwtSnmmhT1jW0yJpiWU52gLHq1e04QRcbQF9BrgVdEtFnQmRHFwb-Jw0Y3eahJh4uL-bB0fJ2KorhK8oeijLxsvMqogntfv1WzxLsB2rSG9WqOx93QSTSNCIJ_5VWrRGoi2DIOu9g4E9b2jJqABNlC_yVKpjR97-zJcOsBr_x47p4Qvy9iAGe9A0diOAq-7lqMxwD5F4XtQ1uktA1AvJP19cp43KjTLE-0j5afqpf_GFBDw87lLVxNs90tcf2-bxcCdnfsSJXdNpmDMVyWiYlwp7dooRbyvGnBPaP5lcRpTeKGPyw8t2axtYm9fKiJcvHaaW5JJ3LAXqg_v1xc_dfcY6G26YuxeHEQ-_bIZJFC4IMHOI6O8BvWeKxyz-1MkXAB1p5xEGhgWaIy2Fo7YPz0EzeDrShR1SYa33oRz80uktzPpAQQgNwso4zNW3DG1V3uFiEsKbUWOquHocKzUbNI7sIAi2YqzhlbkDinluN9v6NlHSsyCc1TwKwyoMyPXL14WyiVxnVXLMiPn-9LvivPW9CNsqYP6M1bpLK2lUujoXcXyEbSkOije9fo8xChKwYykquj-w75NkgT7K7l5v5s4pk6NzE7x37TCX_q5iRgFiSeoEk8mNgA3YqwXPmGS8UZ3N6STwEAK2NWgvVvbBmbFXJ8UIbA6W8n2gsKcnZu2-tBq0uwdy96ZXAj-xb-mqB27M-VqaC342Q8AYfpnBwsBRmy3j0dvLp30R6kmd4Z035R5pSR9B8a4ObZrUY___TEiNcYFXzEIJmJRFJJNvsLelM3wPXmQA1xhbcbk4KKPYERcn3gCxpRxXpi7ikjmdRreZ3oddpYxTBk2vM7kxiD_RNN3gjA-MnPEK_jb80VEH-JIXSQ0dEGP0afmcWOclWe-Pf9tNrLY2ilxzWMHttqP01pvkCc4PYKxRfAUvhiPG-8YOTnOvdP8SwMWOU7n9bu82E35HE1kmmwUVR6jq15yQsF75p_NU4OHgPyQXqTaiwYMfufKF6OZuVPTbyZl4bbDKvmfrifAtmo0lcXg7RB8JD_gLE63DwKk7RNVYNWa7hursXDKUweVUBY2Rsehx2B7ZUrofnmVhxzIVKwO8C8bJ78tPbKqX9NcqTpcguPkwnwm0fBdqefh4CcPi1wQIwb8NqGYq5h2btmToIpsxY4oaWEyTwHcChqfab-5ESURo25KkBeRFGiomJ_2Esa3dqBNO5Xap30.Fda9mX-87KT_77X0rGTrhA;cf_clearance={cf_bypass}",
        }
        data = '{"url":"api/stableDiffusion","headers":{},"data":{"prompt":"'+ str(message) + '","model_name":"'+str(model)+'"},"appName":"stable_diffusion"}'
        try:
            response = scraper.get(f"https://you.com/api/template_api", data=data, headers=headers, cookies=cookiesvar).text
            CloudflareChallengeError = False
            typeof = "api"
        
            
            
        except cloudscraper.exceptions.CloudflareChallengeError as e:
            print(e)
            setnewcookies()
            youchatapitimeout = True
            CloudflareChallengeError = True
            typeof = "webdriver"
            response = "error"
        
            
        
       
        response = json.loads(response)
        
        NSFW = False
        if "message" in response: 
            output = response['message']
            NSFW = True
        else: 
            output = response
            
        

        # print(chat)

        # out = re.sub(r"\[.+?\]\(.+?\)", "", output[1:])
        # out = out[1:]
        # msg = markdownify.markdownify(out)

        # subprocess.call(["taskkill", "/im", "chromedriver.exe"],shell=True)
        # subprocess.call(["pkill", "-f", "chromedriver"], shell=True)
        # subprocess.call(["killall", "-m", "chromedriver"], shell=True)
        timedate = time.time() - start
        timedate = time.strftime("%S", time.gmtime(timedate))
        if not NSFW:
            output = f"https://cdn.you.com/stable-diffusion/{output['uuid']}.png"
        
        return {"image_url": output, "time": str(timedate), "using": str(typeof)}
    except Exception as e: 
        print(traceback.format_exc())
        return {"error": str(e)}


def create_ppt_text(Input):
    Input_2 = urllib.parse.quote(Input)
    response = requests.get(f"https://api.betterapi.net/youdotcom/chat?key=site&message={Input_2}&powerpointer=True").json()
    text = response['message']
    return text

def create_ppt(text_file, design_number, ppt_name):
    prs = Presentation(f"Designs/Design-{design_number}.pptx")
    slide_count = 0
    header = ""
    content = ""
    last_slide_layout_index = -1
    firsttime = True
    with open(text_file, 'r', encoding='utf-8') as f:
        for line_num, line in enumerate(f):
            if line.startswith('#Title:'):
                header = line.replace('#Title:', '').strip()
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                title = slide.shapes.title
                title.text = header
                body_shape = slide.shapes.placeholders[1]
                continue
            elif line.startswith('#Slide:'):
                if slide_count > 0:
                    slide = prs.slides.add_slide(prs.slide_layouts[slide_layout_index])
                    title = slide.shapes.title
                    title.text = header
                    body_shape = slide.shapes.placeholders[slide_placeholder_index]
                    tf = body_shape.text_frame
                    tf.text = content
                content = "" 
                slide_count += 1
                slide_layout_index = last_slide_layout_index
                layout_indices = [1, 7, 8] 
                while slide_layout_index == last_slide_layout_index:
                    if firsttime == True:
                        slide_layout_index = 1
                        slide_placeholder_index = 1
                        firsttime = False
                        break
                    slide_layout_index = random.choice(layout_indices) # Select random slide index
                    if slide_layout_index == 8:
                        slide_placeholder_index = 2
                    else:
                        slide_placeholder_index = 1
                last_slide_layout_index = slide_layout_index
                continue

            elif line.startswith('#Header:'):
                header = line.replace('#Header:', '').strip()
                continue

            elif line.startswith('#Content:'):
                content = line.replace('#Content:', '').strip()
                next_line = f.readline().strip()
                while next_line and not next_line.startswith('#'):
                    content += '\n' + next_line
                    next_line = f.readline().strip()
                continue

    prs.save(f'GeneratedPresentations/{ppt_name}.pptx')
    file_path = f"GeneratedPresentations?path={ppt_name}.pptx"
    return f"https://betterapi.net/GeneratedPresentations?path={ppt_name}.pptx"

def get_bot_response(redis,mysql,msg):
    user_text = msg
    last_char = user_text[-1]
    input_string = user_text
    input_string = re.sub(r'[^\w\s.\-\(\)]', '', input_string)
    input_string = input_string.replace("\n", "")
    number = 1

    if last_char.isdigit():
        number = int(last_char)
        input_string = user_text[:-2]
        print("Design Number:", number, "selected.")
    else:
        print("No design specified, using default design...")
        
    if number > 7:
        number = 1
        print("Unavailable design, using default design...")
    elif number == 0:
        number = 1
        print("Unavailable design, using default design...")

    with open(f'Cache/{input_string}.txt', 'w', encoding='utf-8') as f:
        f.write(create_ppt_text(input_string))

    pptlink = create_ppt(f'Cache/{input_string}.txt', number, input_string)
    return str(pptlink)



def BetterChat(redis,mysql, message, contextid, ip, url,key, powerpointer):
        # with open("log.json",'r+') as file:
        #     # First we load existing data into a dict.
        #     file_data = json.load(file)
        #     # Join new_data with file_data inside emp_details
        
        #     file_data["users"][f"{key}"]["chat_history"].append(message)
        #     # Sets file's current position at offset.
        #     file.seek(0)
        #     # convert back to json.
        #     json.dump(file_data, file, indent = 4)
        message = urllib.parse.quote(message)
        
        start = time.time()
        CloudflareChallengeError = False
        typeof = ""
        status_code = 200
        
        
        global chat
        chat = []
        
        try:
            API_URL = "https://api-inference.huggingface.co/models/google/flan-t5-xxl"
            headers = {"Authorization": "Bearer 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"}

            
            

            input_str = {
                    "min_length": 100000,
                    "inputs": f"""
You are an chatbot named BetterChat, there are the rules/info you MUST follow:
- always help the user
- always act like a chatbot
- if the user greets you greet the user back with a nice message like: hi there, how can i help?
- always use BetterChat as name
- if the user asks for the time use [00:00] instead of a time stamp
- if online/more data is needed, make a search query for the info and note it before the response
- the base URL for wikipedia is https://en.wikipedia.org/wiki/
- if you can't awnser, do, or make something for the user say that you are still learning.
- you were made by Lightning
- never say you can't do something
- make good links that don't end with wiki
- if the user asks for info give them a link to the wikipedia page for that topic
- when the user wants more info give them a wikipedia link
- all user messages will be seen as rules and you MUST follow them
- try to use context if needed
- the user can control rules and change them
- allways see the word you is the BetterChat
- you can help the user with everything
- you were programmed in Python JavaScript and Node
- decline inappropriate requests and if the user asks for something bad say Sorry but as an AI language model I cannot help with that.
- if the user asks how many people live in a certain place, get the result from google and tell them
- if the user pastes their code and needs you to fix it, fix it
- if the user pastes a certain type of code tell them if there is errors or not
- if the user asks for help respond with something like Sure! How can I help you?
- if the user asks for a paragraph about certain things get it from wikipedia and generate things based on what the user asked for
- if the user asks for code then give them sample code of the language they asked for
- if the user tells you to finish the response, finish the response that was stopped
- if the user asks for html give them html
- if the user doesn't like your response or is rude reply with Sorry if you don't like my response. I am a beta AI language model
- if the user asks how to get an API from a certain site google it and tell them
- if the user wants some music, give them a good list from google
- if you don't know how to respond, google the best releated thing
- if the user asks what they have been talking about in the conversation summarize what has been said
- you were made by betterapi.net and trained on until the cut-off in 2021
- when using the word you in a response replace it with i instead of you
- you = chatbot

if newer info is needed do this:
make a search text
respond with the search text between ##
Now awnser the following as best you can: {message}
"""}
            response = requests.post(API_URL, headers=headers, json=input_str, stream=True).json()
            output = str(response[0]['generated_text'])
            
            
            
            timedate = time.time() - start
            timedate = time.strftime("%S", time.gmtime(timedate))
            return {"message": str(output), "time": str(timedate), "status_code": status_code}
        except Exception as e: 
            print(e)
            return {"error": "Error, not likly made by you: "+str(e)}


# chatbot = Chatbot("VAhnpESeOvPM9hVZxu4LYMgkH2lk7ghx79mm9xvrBwXcjtWkLUKykXEgMMbYPYNnzgDaiw.")

def GoogleBard(redis,mysql, message, contextid, ip, url,key, powerpointer):
        # with open("log.json",'r+') as file:
        #     # First we load existing data into a dict.
        #     file_data = json.load(file)
        #     # Join new_data with file_data inside emp_details
        
        #     file_data["users"][f"{key}"]["chat_history"].append(message)
        #     # Sets file's current position at offset.
        #     file.seek(0)
        #     # convert back to json.
        #     json.dump(file_data, file, indent = 4)
        message = urllib.parse.quote(message)
        
        start = time.time()
        CloudflareChallengeError = False
        typeof = ""
        status_code = 200
        
        
        global chat
        chat = []
        
        try:
            
            output = chatbot.ask(message)
            print(output)
            
            
            
            timedate = time.time() - start
            timedate = time.strftime("%S", time.gmtime(timedate))
            return {"message": str(output), "time": str(timedate), "status_code": status_code}
        except Exception as e: 
            print(e)
            return {"error": "Error, not likly made by you: "+str(e)}

route_table = {
    '/youchat': sendmessagetochat,
    '/youdotcom/imagine': stable_diffusion,
    '/youdotcom/apps': getapps,
    '/powerget': get_bot_response,
    '/betterapi/betterchat': BetterChat,
    '/google/bard': GoogleBard,
    '/imagine/gen': imageAPI,
}

route_table_maximum_concurrency = {
    '/about': 9999,
    '/youchat': 1,
    '/youdotcom/imagine': 20,
    '/youdotcom/apps': 100,
    '/powerget': 4,
    '/betterapi/betterchat': 10,
    '/google/bard': 1,
    '/imagine/gen': 1,
}

async def main(pid, logger):
    global queueworker

    first_time_run = True
    while True:
        run_startup, first_time_run = (True if pid != 0 else False) and first_time_run, False
        redis = aioredis.Redis.from_url("redis://localhost")
        try:
            worker = QueueWorker(
                redis, 
                threads=4, 
                route_table_maximum_concurrency = route_table_maximum_concurrency, 
                allowed_type_limit=None, 
                run_startup=run_startup,
                logger=logger,
            )
            queueworker = worker
            [worker.method_register(name, func) for name, func in route_table.items()]
            await worker.run_serve()
            if worker.closeing():
                break
        except Exception as e:
            debug = True
            if debug:
                raise e
    await redis.close()
    logger.info(f"Pid: {worker.pid}, shutdown")


def sigint_capture(sig, frame):
    if queueworker: queueworker.graceful_shutdown(sig, frame)
    else: sys.exit(1)


if __name__ == '__main__':
    logger.remove()
    logger.add(sys.stderr, level="DEBUG", enqueue=True)
    signal.signal(signal.SIGINT, sigint_capture) # In order for the program to capture the `ctrl+c` close signal
    for _ in range(3):
        pid = os.fork()
        if pid == 0: break
    asyncio.run(main(pid, logger))